"""
This script is designed to extract text from various document formats (.pdf, .docx, .doc),
chunk the text into smaller sections, and embed these sections, and save them to a vector store.

This is part of a larger pipeline for document processing and embedding, and is run upon the upload of new documents.

Supported file formats:
- PDF
- DOCX, DOC
- XLSX, XLS
- PPTX, PPT
- CSV
- TXT
- HTML
- Markdown

MIT License
Copyright (c) 2025 Tey Xue Cong, Tan Hong Kai, Siah Wan Ru Tricia, Tee Jia Yee
See the LICENSE file in the project root for full license information.
"""
import warnings
warnings.filterwarnings("ignore", category=RuntimeWarning, message=".*grpcio.*", module="opentelemetry.*")
import io
import json
import unicodedata
import xlrd
from openpyxl import load_workbook
from pptx import Presentation
from docx import Document
import json
import pdfplumber
from lxml import etree
from langchain.text_splitter import RecursiveCharacterTextSplitter
import pandas as pd
import sys
import win32com.client
from pathlib import Path
from operator import itemgetter
import chromadb
from typing import List
from typing import Union
from sentence_transformers import SentenceTransformer
from llama_index.core.schema import Document as llamadoc
from llama_index.core import VectorStoreIndex, StorageContext
from llama_index.embeddings.huggingface import HuggingFaceEmbedding
from llama_index.core import Settings
from llama_index.llms.ollama import Ollama
from llama_index.core.node_parser import SimpleNodeParser
from bs4 import BeautifulSoup, Tag, NavigableString
from llama_index.vector_stores.chroma import ChromaVectorStore
from llama_index.core.base.embeddings.base import BaseEmbedding


# Use pathlib so it works on both Windows and Linux
PROJECT_ROOT = Path(__file__).resolve().parent.parent
print(f"Project root: {PROJECT_ROOT}")
RAW_DATA = PROJECT_ROOT / 'data' / 'raw_data'
LOG_FILE = PROJECT_ROOT / 'data' / 'Logs' / 'processed_files.json'
PERSIST_DIR = PROJECT_ROOT / 'data' / 'ChromaDB'

# This one is for embeding USER query
Settings.embed_model = HuggingFaceEmbedding(model_name="intfloat/e5-large-v2") # MUST BE SAME AS THE ONE USED FOR INDEXING
# Load model locally
model = SentenceTransformer("intfloat/e5-large-v2")

# ChromaDB embedding function
class HFChromaEmbedding:
    def __init__(self, model):
        self.model = model

    def __call__(self, input: Union[str, List[str]]) -> Union[List[float], List[List[float]]]:
        if isinstance(input, str):
            input = [f"passage: {input}"]
            return self.model.encode(input, convert_to_numpy=True).tolist()[0]
        else:
            input = [f"passage: {text}" for text in input]
            return self.model.encode(input, convert_to_numpy=True).tolist()

    def name(self) -> str:
        return "HFChromaEmbedding-e5-large-v2"

chroma_embedding_fn = HFChromaEmbedding(model)

# LlamaIndex embedding class
class HFLlamaEmbedding(BaseEmbedding):
    model: SentenceTransformer
    def __init__(self, model):
        super().__init__(model=model)

    def _get_text_embedding(self, text: str) -> List[float]:
        return self.model.encode(f"passage: {text}", convert_to_numpy=True).tolist()

    def _get_text_embeddings(self, texts: List[str]) -> List[List[float]]:
        return self.model.encode([f"passage: {t}" for t in texts], convert_to_numpy=True).tolist()

    def _get_query_embedding(self, query: str) -> List[float]:
        return self.model.encode(f"query: {query}", convert_to_numpy=True).tolist()

    async def _aget_query_embedding(self, query: str) -> List[float]:
        return self._get_query_embedding(query)

llama_embedding_fn = HFLlamaEmbedding(model)

# The embedding model is chosen based on the Hugging Face MTEB leaderboard:
# https://huggingface.co/spaces/mteb/leaderboard?benchmark_name=MTEB(Multilingual,+v2)

# This model is a subset of the multilingual model "multilingual-e5-large-instruct", which was 4th on the leaderboard.
# We selected "intfloat/e5-large-v2" as it offers the best performance-to-efficiency ratio
# It is faster and more lightweight than the multilingual version, making it suitable for our real-time use case.

# Set Ollama as the default LLM globally
Settings.llm = Ollama(model="llama3.2", context_window=4096, timeout=120)

faiss_path = "faiss.index"
faiss_file_path = PERSIST_DIR / faiss_path

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

# Departments and countries for metadata
ALL_DEPARTMENTS = [
    "Human Resource", "Admin & Operations", "Project Management", "Procurement",
    "IT", "Marketing", "Business Development", "Finance", "Service Delivery"
]
ALL_COUNTRIES = [
    "Singapore", "United Kingdom", "United States", "Thailand", "Indonesia",
    "Korea", "China", "Japan", "Vietnam", "Myanmar"
]


# 1. Extraction

def extract_text_from_pdf(file_path):
    """Extracts .pdf files"""
    output = []

    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            blocks = []

            # Extract all tables with bbox
            tables = page.find_tables()
            for table in tables:
                table_bbox = table.bbox
                table_content = []
                for row in table.extract():
                    row_text = " | ".join(cell.strip() if cell else "" for cell in row)
                    table_content.append(f"| {row_text} |")
                blocks.append({
                    'type': 'table',
                    'top': table_bbox[1],
                    'bottom': table_bbox[3],
                    'content': "\n".join(table_content)
                })

            # Extract all words
            words = page.extract_words()
            # Group words into lines by their vertical position (rounded)
            lines_map = {}
            for word in words:
                top = round(word['top'], 1)
                if top not in lines_map:
                    lines_map[top] = []
                lines_map[top].append(word)

            # Convert lines_map to list of text blocks
            for top, word_group in lines_map.items():
                line_text = " ".join(w['text'] for w in sorted(word_group, key=lambda w: w['x0']))
                # Check if line overlaps any table
                in_table = False
                for t in blocks:
                    if t['type'] == 'table' and t['top'] <= top <= t['bottom']:
                        in_table = True
                        break
                if not in_table:
                    blocks.append({
                        'type': 'text',
                        'top': top,
                        'content': line_text
                    })

            # Sort blocks by Y position
            blocks_sorted = sorted(blocks, key=itemgetter('top'))

            page_output = [block['content'] for block in blocks_sorted]
            output.append("\n".join(page_output))

    return "\n\n".join(output)


def extract_text_from_doc(file_path):
    """Extracts .doc files, ignores images"""

    # Have to open the file in background because its old -_-
    word = win32com.client.Dispatch("Word.Application")
    word.Visible = False
    doc = word.Documents.Open(file_path)

    full_text = []
    content = doc.Content
    start = content.Start
    end = content.End
    bullets = {'•', '‣', '·', '‧', '–', '-', '*', '', '●', '■', '♦', '\uf0b7', 'o'}

    while start < end:
        current_range = doc.Range(start, start + 1)
        # Extract tables
        if current_range.Tables.Count > 0:
            table = current_range.Tables(1)
            table_text = []
            for row in table.Rows:
                row_text = []
                for cell in row.Cells:
                    cell_text = cell.Range.Text.strip().replace('\r', '').replace('\x07', '')
                    row_text.append(cell_text)
                table_text.append(' | '.join(row_text))
            full_text.append('\n'.join(table_text))
            start = table.Range.End

        # Extract paragraphs
        elif current_range.Paragraphs.Count > 0:
            para_range = current_range.Paragraphs(1).Range
            para_text = para_range.Text.strip().replace('\r', '').replace('\x07', '')
            para_text = unicodedata.normalize("NFKC", para_text) # Normalize to raw text

            if para_text:
                list_format = para_range.ListFormat
                indent = ''

                '''This chunk of code is to deal with microsoft word lists'''
                if list_format.ListType != 0:
                    # Get list indent level and marker
                    level = max(list_format.ListLevelNumber, 1)
                    indent = '    ' * (level - 1)
                    marker = list_format.ListString.strip()

                    # Some markers are invisible or from Wingdings/Symbol font (like '\uf0b7')
                    # These don't render well, so we substitute a standard bullet
                    if not marker or not marker.isprintable() or ord(marker[0]) >= 0xF000:
                        marker = '•'
                    para_text = f"{indent}{marker} {para_text}"

                else:
                    stripped = para_text.lstrip()
                    # Now check if the first character is a bullet point
                    if stripped and stripped[0] in bullets:
                        para_text = f"• {stripped[1:].lstrip()}"

                full_text.append(para_text)

            start = para_range.End
        else:
            start = current_range.End  # Fallback to avoid infinite loop

     # Extract text from shapes
    for shape in doc.Shapes:
        if shape.TextFrame.HasText:
            shape_text = shape.TextFrame.TextRange.Text.strip().replace('\r', '').replace('\x07', '')
            if shape_text:
                full_text.append("[Shape Text] " + shape_text)

    for ishape in doc.InlineShapes:
        if hasattr(ishape, "TextFrame") and ishape.TextFrame.HasText:
            shape_text = ishape.TextFrame.TextRange.Text.strip().replace('\r', '').replace('\x07', '')
            if shape_text:
                full_text.append("[Inline Shape Text] " + shape_text)
                
    doc.Close(False)
    word.Quit()
    return '\n'.join(full_text)


def extract_text_from_docx(file_path):
    """Extracts .docx files, ignores images"""
    doc = Document(file_path)
    full_text = []

    # Get the raw XML of the Word doc so we can manually look at paragraphs, tables, etc.
    doc_xml = doc.element
    namespaces = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}

    for element in doc_xml.body:
        # Extract text
        if element.tag == etree.QName(namespaces['w'], 'p'):
            # Detect if it's a list item
            num_pr = element.find('.//w:numPr', namespaces)
            is_list = num_pr is not None

            # Extract hyperlink stuff
            hyperlink = element.find('.//w:hyperlink', namespaces)
            if hyperlink is not None:

                # Hyperlink text
                texts = hyperlink.findall('.//w:t', namespaces)
                link_text = ''.join(t.text for t in texts if t.text)

                # Get the hyperlink target
                r_id = hyperlink.attrib.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                if r_id:
                    rels = doc.part.rels
                    url = rels[r_id]._target if r_id in rels else ''
                    plain_text = f"[{link_text}]({url})"
                else:
                    # No url found
                    plain_text = link_text
            else:
                # Plain text
                texts = element.findall('.//w:t', namespaces)
                plain_text = ''.join(t.text for t in texts if t.text)

            plain_text = plain_text.strip()
            if plain_text:
                if is_list:
                    plain_text = f"• {plain_text}"
                full_text.append(plain_text)

        # Extract tables
        elif element.tag == etree.QName(namespaces['w'], 'tbl'):
            for row in element.findall('.//w:tr', namespaces):
                cells = row.findall('.//w:tc', namespaces)
                row_cells = []
                for cell in cells:
                    texts = cell.findall('.//w:t', namespaces)
                    cell_text = ''.join(t.text for t in texts if t.text).strip()
                    row_cells.append(cell_text)
                formatted_row = '| ' + ' | '.join(row_cells) + ' |'
                full_text.append(formatted_row)

    return '\n'.join(full_text)


def extract_text_from_xls(file_path):
    """Extracts .xls file, ignoring formulas"""
    book = xlrd.open_workbook(file_path)
    all_text = ""

    for sheet in book.sheets():
        all_text += f"--- Sheet: {sheet.name} ---\n"
        for row_idx in range(sheet.nrows):
            row_values = sheet.row_values(row_idx)
            line_parts = []
            for val in row_values:
                if isinstance(val, float):
                    line_parts.append(f"{val:.7g}")  # Prevent floating points stuff from going crazy
                else:
                    line_parts.append(str(val).strip())
            all_text += "\t".join(line_parts) + "\n"
        all_text += "\n"
    
    return all_text


def extract_text_from_xlsx(file_path):
    """Extracts .xlsx files, ignoring formulas"""
    wb = load_workbook(filename=file_path, data_only=True) # Don't extract formulas
    all_text = ""

    for sheet in wb.sheetnames:
        ws = wb[sheet]
        all_text += f"--- Sheet: {sheet} ---\n"
        for row in ws.iter_rows():
            row_values = []
            for cell in row:
                value = str(cell.value) if cell.value is not None else ""
                row_values.append(value)
            all_text += "\t".join(row_values).rstrip() + "\n"
        all_text += "\n"
    
    return all_text


def extract_text_from_csv(file_path):
    """Extracts .csv files"""
    df = pd.read_csv(file_path)
    return df.to_string(index=False)


def extract_text_from_txt(file_path):
    """Extracts .txt files"""
    with open(file_path, "r", encoding="utf-8") as file:
        return file.read()
    

def extract_text_from_html(file_path):
    """Extract .html files"""
    with open(file_path, 'r', encoding='utf-8') as file:
        soup = BeautifulSoup(file, 'html.parser')

    def format_element(element):
        """
        Apply custom formatting to certain HTML elements:
        - <li>: Format as list item with a dash
        - <table>: Format as a tab-separated grid
        - <a>: Replace with "text (href)" format
        """
        if element.name == 'li':
            return f"- {element.get_text(' ', strip=True)}\n"

        elif element.name == 'table':
            table_text = []
            for row in element.find_all('tr'):
                row_text = []
                for cell in row.find_all(['td', 'th']):
                    # Handles cells in tables
                    cell_text = cell.get_text(" ", strip=True).replace('\r', '').replace('\x07', '')
                    row_text.append(cell_text)
                table_text.append(' | '.join(row_text))
            return '\n'.join(table_text) + '\n'
        
        elif element.name == 'a' and element.has_attr('href'):
            return f"{element.get_text(strip=True)} ({element['href']})"
        else:
            return None

    def traverse(node):
        """Traverse the HTML tree and extract formatted text."""
        output = ''

        for child in node.children:
            if isinstance(child, NavigableString):
                output += child
            elif isinstance(child, Tag):
                if child.name in ['script', 'style']:
                    continue
                # Deal with hyperlinks
                if child.name == 'a' and child.has_attr('href'):
                    output += format_element(child)
                elif child.name in ['li', 'table']:
                    output += format_element(child)
                else:
                    output += traverse(child)
        return output

    body = soup.body or soup  # Fallback if the text is not associated with a <body> thing
    formatted_text = traverse(body)
    return "\n".join(line.strip() for line in formatted_text.splitlines() if line.strip())


def extract_text_from_ppt(ppt_path):
    """Extract text and speaker notes from .ppt"""
    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
    powerpoint.Visible = 1

    presentation = powerpoint.Presentations.Open(ppt_path, WithWindow=False)
    all_text = []

    for i, slide in enumerate(presentation.Slides, start=1):
        slide_text = [f"--- Slide {i} ---"]

        for shape in slide.Shapes:
            # Handle text and bullet lists
            if shape.HasTextFrame:
                tf = shape.TextFrame
                if tf.HasText:
                    paragraphs = []
                    for paragraph in tf.TextRange.Paragraphs():
                        text = paragraph.Text.strip().replace('\r', '')
                        if text:
                            bullet = "- " if paragraph.ParagraphFormat.Bullet.Type != 0 else ""
                            paragraphs.append(bullet + text)
                    if paragraphs:
                        slide_text.append("\n".join(paragraphs))

            # Handle tables
            if shape.HasTable:
                table = shape.Table
                table_text = []
                for row in range(1, table.Rows.Count + 1):
                    row_text = []
                    for col in range(1, table.Columns.Count + 1):
                        cell = table.Cell(row, col)
                        cell_text = cell.Shape.TextFrame.TextRange.Text.strip().replace('\r', '').replace('\x07', '')
                        row_text.append(cell_text)
                    table_text.append(" | ".join(row_text))
                slide_text.append("\n".join(table_text))

        # Speaker Notes
        if slide.NotesPage.Shapes.Placeholders.Count >= 2:
            notes_shape = slide.NotesPage.Shapes.Placeholders(2)
            if notes_shape.HasTextFrame and notes_shape.TextFrame.HasText:
                notes = notes_shape.TextFrame.TextRange.Text.strip().replace('\r', '')
                if notes:
                    slide_text.append(f"[Notes] {notes}")

        all_text.append("\n".join(slide_text))

    presentation.Close()
    powerpoint.Quit()

    return "\n\n".join(all_text)


def is_bullet_paragraph(paragraph):
    """
    Check if a paragraph has bullet formatting by inspecting XML.
    We have to do this since python-pptx cannot detect bullets natively :(
    """
    pPr = paragraph._element.pPr
    return pPr is not None and pPr.find(".//a:buChar", namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}) is not None

def extract_text_from_pptx(file_path):
    """Extract text and notes from a .pptx file."""
    presentation = Presentation(file_path)
    all_text = []

    for i, slide in enumerate(presentation.slides, start=1):
        slide_text = [f"--- Slide {i} ---"]

        for shape in slide.shapes:
            # Handle text and bullet lists
            if shape.has_text_frame:
                paragraphs = []
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue

                    if is_bullet_paragraph(para):
                        indent = "  " * para.level
                        paragraphs.append(f"{indent}- {text}")
                    else:
                        paragraphs.append(text)

                if paragraphs:
                    slide_text.append("\n".join(paragraphs))

            # Handle Tables
            if shape.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE
                table = shape.table
                table_text = []
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        cell_text = cell.text.strip().replace('\r', '').replace('\x07', '')
                        row_text.append(cell_text)
                    table_text.append(' | '.join(row_text))
                slide_text.append('\n'.join(table_text))

        # Speaker Notes
        notes_slide = slide.notes_slide if slide.has_notes_slide else None
        if notes_slide:
            notes_text = notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_text.append(f"[Notes] {notes_text}")

        all_text.append("\n\n".join(slide_text))

    return "\n\n".join(all_text)

def extract_text_from_md(file_path):
    """Extract .md file"""
    with open(file_path, 'r', encoding='utf-8') as file:
        return file.read()


# 2. Chunking

def split_into_documents(text, chunk_size=1000, chunk_overlap=200, title="Untitled", source="unknown.txt", departments=None, countries=None):
    """
    Splits text into chunks and returns them as LlamaIndex Document objects with metadata.

    Parameters:
        text (str): The full input text to split.
        chunk_size (int): Max characters per chunk.
        chunk_overlap (int): Characters to overlap between chunks.
        title (str): Title of the source document.
        source (str): Filename of the source document.
        departments (list): List of departments that can see the document.
        countries (list): List of countries that can see the document.

    Returns:
        List[Document]: Chunked Document objects with metadata.
    """
    if departments is None:
            departments = []
    if countries is None:
        countries = []

    splitter = RecursiveCharacterTextSplitter( # Used RecussiveCharacterTextSplitter because it's good at identifying paragraphs and natural sections
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap
    )

    chunks = splitter.split_text(text)
    print(f"Total chunks: {len(chunks)}\n", flush=True)

    documents = []
    for i, chunk in enumerate(chunks):
        metadata = {
            "chunk": i,
            "title": title,
            "source": source,
        }
        for dept in ALL_DEPARTMENTS:
            metadata[dept] = str(dept in departments)
        for country in ALL_COUNTRIES:
            metadata[country] = str(country in countries)

        doc = llamadoc(text=chunk, metadata=metadata)
        documents.append(doc)

    return documents

# 3. Embedding + Saving

def build_or_append_index(documents,embed_model,persist_dir="pipeline/data/EmbeddedChroma",collection_name="quickstart"):
    """
    Create or append to a ChromaDB + LlamaIndex index.

    Parameters:
        documents (List[Document]): New documents to insert
        embed_model (BaseEmbedding): Embedding model (e.g., HuggingFaceEmbedding)
        persist_dir (str or Path): Directory where ChromaDB data is stored
        collection_name (str): Name of the ChromaDB collection
    """
    persist_dir = Path(persist_dir)
    persist_dir.mkdir(parents=True, exist_ok=True)

    # Initialize ChromaDB persistent client
    db = chromadb.PersistentClient(path=str(persist_dir))
    chroma_collection = db.get_or_create_collection(collection_name, embedding_function=chroma_embedding_fn)
    vector_store = ChromaVectorStore(chroma_collection=chroma_collection)
    storage_context = StorageContext.from_defaults(vector_store=vector_store)
    

    # Check if collection already has data
    existing_count = chroma_collection.count()
    if existing_count > 0:
        print("Loading existing ChromaDB index...", flush=True)
        index = VectorStoreIndex.from_vector_store(
            vector_store, storage_context=storage_context, embed_model=embed_model
        )
        print("Total docs in index before append:", existing_count, flush=True)
        # Append new documents
        parser = SimpleNodeParser()
        nodes = parser.get_nodes_from_documents(documents)
        index.insert_nodes(nodes)
    else:
        print("Creating new ChromaDB index...", flush=True)
        index = VectorStoreIndex.from_documents(
            documents, storage_context=storage_context, embed_model=embed_model
        )

    # ChromaDB persists automatically, but you can force a flush if needed
    print("Saving new data...", flush=True)
    print("Total docs in index:", chroma_collection.count(), flush=True)

    return index


# 4. Main Pipeline
def extract_text_from_file(file_path):
    """Extract text based on file extension"""
    file_extension = file_path.lower().split('.')[-1]
    
    if file_extension == "docx":
        return extract_text_from_docx(file_path)
    elif file_extension == "doc":
        return extract_text_from_doc(file_path)
    elif file_extension == "pdf":
        return extract_text_from_pdf(file_path)
    elif file_extension == "pptx":
        return extract_text_from_pptx(file_path)
    elif file_extension == "ppt":
        return extract_text_from_ppt(file_path)
    elif file_extension == "xls":
        return extract_text_from_xls(file_path)
    elif file_extension == "xlsx":
        return extract_text_from_xlsx(file_path)
    elif file_extension == "csv":
        return extract_text_from_csv(file_path)
    elif file_extension == "txt":
        return extract_text_from_txt(file_path)
    elif file_extension == "html":
        return extract_text_from_html(file_path)
    elif file_extension == "md":
        return extract_text_from_md(file_path)
    elif file_extension == "json":
        return  # Skip JSON files, as they are used for metadata and not for text extraction
    else:
        raise ValueError(f"Unsupported file type: {file_extension}")


def run_pipeline(raw_folder):
    new_files = []

    log_file_path = Path(LOG_FILE)
    raw_folder = Path(raw_folder)

    log_file_path.parent.mkdir(parents=True, exist_ok=True)

    # Ensure log file exists and has valid JSON
    if not log_file_path.exists() or log_file_path.stat().st_size == 0:
        with open(log_file_path, "w") as f:
            json.dump([], f)

    with open(log_file_path, "r") as f:
        processed_files = set(json.load(f))

    for filename in raw_folder.iterdir():
        if not filename.is_file():
            continue
        if filename.name in processed_files:
            print(f'Skipping file: {filename.name}', flush=True)
            continue  # skip already processed files
        if filename.suffix == ".json" and filename.name.endswith(".meta.json"):
            print(f"Skipping metadata file: {filename.name}")
            continue  # skip metadata sidecar files
        if filename.name.startswith(".batch_upload_") and filename.suffix == ".json":
            print(f"Skipping batch tracking file: {filename.name}")
            continue # skip batch tracking files

        print(f"\nProcessing file: {filename.name}", flush=True)
        extracted_text = extract_text_from_file(str(filename))

        meta_path = filename.with_suffix(filename.suffix + ".meta.json")
        departments = []
        countries = []
        if meta_path.exists():
            with open(meta_path, "r", encoding="utf-8") as f:
                meta = json.load(f)
                departments = meta.get("departments", [])
                countries = meta.get("countries", [])


        documents = split_into_documents(
        extracted_text,
        title=filename.name,
        source=filename.name,
        departments=departments,
        countries=countries
    )
        
        build_or_append_index(documents, llama_embedding_fn, persist_dir=PERSIST_DIR, collection_name="quickstart")
        new_files.append(filename.name)

    # Update log file
    processed_files.update(new_files)
    with open(log_file_path, "w") as f:
        json.dump(list(processed_files), f, indent=2)

    print("Pipeline completed. Database has been updated!", flush=True)

if __name__ == "__main__":
    run_pipeline(RAW_DATA)
